import csv
import json
from io import BytesIO, StringIO
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas

from backend.app.main import create_app


def make_docx(text: str) -> bytes:
    stream = BytesIO()
    document = Document()
    document.add_heading("StudyPilot", level=1)
    document.add_paragraph(text)
    document.save(stream)
    return stream.getvalue()


def make_pdf(text: str) -> bytes:
    stream = BytesIO()
    page = canvas.Canvas(stream)
    page.drawString(72, 760, text)
    page.save()
    return stream.getvalue()


def make_blank_pdf() -> bytes:
    stream = BytesIO()
    page = canvas.Canvas(stream)
    page.showPage()
    page.save()
    return stream.getvalue()


def make_notebook() -> bytes:
    return json.dumps(
        {
            "nbformat": 4,
            "nbformat_minor": 5,
            "metadata": {"kernelspec": {"name": "python3", "display_name": "Python 3"}},
            "cells": [
                {
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": ["# Experiment\n", "Explain the **result**."],
                },
                {
                    "cell_type": "code",
                    "execution_count": 7,
                    "metadata": {},
                    "source": ["score = 0.95\n", "score"],
                    "outputs": [
                        {
                            "output_type": "execute_result",
                            "execution_count": 7,
                            "metadata": {},
                            "data": {"text/plain": ["0.95"]},
                        }
                    ],
                },
            ],
        },
        ensure_ascii=False,
    ).encode("utf-8")


def make_styled_xlsx() -> bytes:
    stream = BytesIO()
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Styled"
    worksheet.merge_cells("A1:B1")
    cell = worksheet["A1"]
    cell.value = "Quarterly report"
    cell.font = Font(name="Aptos", size=16, bold=True, color="FFFFFF")
    cell.fill = PatternFill(fill_type="solid", fgColor="2F5597")
    cell.alignment = Alignment(horizontal="center", vertical="center")
    worksheet.column_dimensions["A"].width = 24
    worksheet.row_dimensions[1].height = 30
    workbook.save(stream)
    return stream.getvalue()


def make_xlsx(text: str) -> bytes:
    stream = BytesIO()
    with ZipFile(stream, "w", ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
              <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "xl/workbook.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
              <sheets><sheet name="概念表" sheetId="1" r:id="rId1"/></sheets>
            </workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
            <worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
              <sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>{text}</t></is></c></row></sheetData>
            </worksheet>""",
        )
    return stream.getvalue()


def make_pptx(text: str) -> bytes:
    stream = BytesIO()
    with ZipFile(stream, "w", ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
              <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "ppt/presentation.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
              <p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>
              <p:sldSz cx="9144000" cy="6858000" type="screen4x3"/>
              <p:notesSz cx="6858000" cy="9144000"/>
            </p:presentation>""",
        )
        archive.writestr(
            "ppt/_rels/presentation.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "ppt/slides/slide1.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
            <p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
              <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>
                <p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{text}</a:t></a:r></a:p></p:txBody></p:sp>
              </p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
            </p:sld>""",
        )
    return stream.getvalue()


def make_styled_pptx() -> bytes:
    stream = BytesIO()
    presentation = PptxPresentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1.2))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Positioned title"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2F, 0x55, 0x97)
    presentation.save(stream)
    return stream.getvalue()


def test_import_preserves_source_creation_and_import_times(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        response = client.post(
            "/api/documents/import",
            files={"file": ("timeline.md", b"# Timeline", "text/markdown")},
            data={"source_created_at": "2025-02-03T04:05:06.000Z"},
        )
        assert response.status_code == 201, response.text
        imported = response.json()["data"]
        listed = client.get("/api/documents").json()["data"][0]

    assert imported["source_created_at"] == "2025-02-03T04:05:06.000Z"
    assert imported["created_at"]
    assert listed["source_created_at"] == imported["source_created_at"]
    assert listed["created_at"] == imported["created_at"]


def test_imports_real_pdf_docx_markdown_and_text(tmp_path) -> None:
    files = [
        ("rag.pdf", make_pdf("BM25 retrieval baseline"), "application/pdf", "BM25 retrieval baseline"),
        (
            "agent.docx",
            make_docx("Tool calling needs idempotency"),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "Tool calling needs idempotency",
        ),
        ("notes.md", "# RRF\n融合稀疏与稠密召回".encode(), "text/markdown", "融合稀疏与稠密召回"),
        ("plain.txt", "Python 子进程隔离".encode(), "text/plain", "Python 子进程隔离"),
    ]

    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = []
        for name, content, media_type, expected in files:
            response = client.post("/api/documents/import", files={"file": (name, content, media_type)})
            assert response.status_code == 201, response.text
            assert expected in response.json()["data"]["body"]
            imported.append(response.json()["data"])

    assert len(imported) == 4
    assert all((tmp_path / item["stored_path"]).is_file() for item in imported)

def test_keeps_unparseable_pdf_as_a_readable_original_file(tmp_path) -> None:
    malformed_pdf = b"%PDF-1.7\nnot a parseable PDF"
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("scanned.pdf", malformed_pdf, "application/pdf")},
        )
        assert imported.status_code == 201, imported.text
        item = imported.json()["data"]
        content = client.get(f"/api/documents/{item['id']}/content")
        original = client.get(f"/api/documents/{item['id']}/file")

    assert item["format"] == "pdf"
    assert item["metadata"]["parse_fallback"] is True
    assert content.status_code == 200
    assert content.json()["data"]["blocks"]
    assert original.status_code == 200

def test_imports_six_formats_as_stable_structured_blocks(tmp_path) -> None:
    cases = [
        ("rag.pdf", make_pdf("BM25 retrieval baseline"), "application/pdf", "pdf", "page"),
        (
            "agent.docx",
            make_docx("Tool calling needs idempotency"),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
            "paragraph",
        ),
        ("notes.md", b"# RRF\nSparse and dense", "text/markdown", "markdown", "section"),
        ("plain.txt", b"Python process isolation", "text/plain", "text", "paragraph"),
        (
            "concepts.xlsx",
            make_xlsx("Feature matrix"),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "xlsx",
            "sheet",
        ),
        ("metrics.csv", "name,score\n\"RRF, hybrid\",0.95\n".encode(), "text/csv", "csv", "sheet"),
        (
            "lecture.pptx",
            make_pptx("Attention mechanism"),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
            "slide",
        ),
    ]

    with TestClient(create_app(data_dir=tmp_path)) as client:
        for filename, content, media_type, expected_format, locator_key in cases:
            response = client.post(
                "/api/documents/import",
                files={"file": (filename, content, media_type)},
            )
            assert response.status_code == 201, response.text
            item = response.json()["data"]
            assert item["format"] == expected_format
            detail = client.get(f"/api/documents/{item['id']}/content")
            assert detail.status_code == 200, detail.text
            blocks = detail.json()["data"]["blocks"]
            assert blocks
            assert blocks[0]["block_key"]
            assert locator_key in blocks[0]["locator"]


def test_csv_import_preserves_quoted_values_and_exposes_editable_cells(tmp_path) -> None:
    payload = 'name,notes,score\n"RRF, hybrid","line 1\nline 2",0.95\n'.encode("utf-8")
    with TestClient(create_app(data_dir=tmp_path)) as client:
        response = client.post(
            "/api/documents/import",
            files={"file": ("metrics.csv", payload, "text/csv")},
        )
        assert response.status_code == 201, response.text
        item = response.json()["data"]
        content = client.get(f"/api/documents/{item['id']}/content").json()["data"]

    assert item["format"] == "csv"
    assert content["blocks"][0]["locator"] == {"sheet": "CSV", "range": "A1:C2"}
    cells = {cell["address"]: cell["value"] for cell in content["blocks"][0]["data"]["cells"]}
    assert cells == {"A1": "name", "B1": "notes", "C1": "score", "A2": "RRF, hybrid", "B2": "line 1\nline 2", "C2": "0.95"}

def test_imports_common_editable_text_and_code_formats(tmp_path) -> None:
    cases = [
        ("config.json", b'{"theme":"glass"}', "application/json"),
        ("pipeline.yaml", b"steps:\n  - import\n", "application/yaml"),
        ("settings.yml", b"enabled: true\n", "application/yaml"),
        ("document.xml", b"<root>StudyPilot</root>\n", "application/xml"),
        ("page.html", b"<h1>StudyPilot</h1>\n", "text/html"),
        ("theme.css", b"body { color: black; }\n", "text/css"),
        ("app.js", b"console.log('StudyPilot');\n", "text/javascript"),
        ("component.jsx", b"export default () => <main />;\n", "text/jsx"),
        ("types.ts", b"export type Mode = 'study';\n", "text/typescript"),
        ("view.tsx", b"export const View = () => <main />;\n", "text/tsx"),
        ("analysis.py", b"print('StudyPilot')\n", "text/x-python"),
        ("Main.java", b"class Main {}\n", "text/x-java-source"),
        ("core.c", b"int main(void) { return 0; }\n", "text/x-c"),
        ("core.cpp", b"int main() { return 0; }\n", "text/x-c++"),
        ("core.h", b"#define STUDY 1\n", "text/x-c"),
        ("core.hpp", b"#pragma once\n", "text/x-c++"),
        ("query.sql", b"select * from notes;\n", "application/sql"),
        ("service.log", b"ready\n", "text/plain"),
        ("app.ini", b"[study]\nenabled=true\n", "text/plain"),
        ("project.toml", b"title='StudyPilot'\n", "application/toml"),
    ]
    with TestClient(create_app(data_dir=tmp_path)) as client:
        for filename, payload, media_type in cases:
            response = client.post(
                "/api/documents/import",
                files={"file": (filename, payload, media_type)},
            )
            assert response.status_code == 201, response.text
            item = response.json()["data"]
            assert item["format"] == "text"
            assert client.get(f"/api/documents/{item['id']}/content").json()["data"]["blocks"]


def test_markdown_and_text_blocks_preserve_source_line_ranges(tmp_path) -> None:
    markdown = "# 第一节\n第一行\n第二行\n\n## 第二节\n结论\n"
    plain = "甲行\n乙行\n\n丙行\n"

    with TestClient(create_app(data_dir=tmp_path)) as client:
        markdown_item = client.post(
            "/api/documents/import",
            files={"file": ("lines.md", markdown.encode("utf-8"), "text/markdown")},
        ).json()["data"]
        plain_item = client.post(
            "/api/documents/import",
            files={"file": ("lines.txt", plain.encode("utf-8"), "text/plain")},
        ).json()["data"]
        markdown_blocks = client.get(
            f"/api/documents/{markdown_item['id']}/content"
        ).json()["data"]["blocks"]
        plain_blocks = client.get(
            f"/api/documents/{plain_item['id']}/content"
        ).json()["data"]["blocks"]

    assert markdown_blocks[0]["locator"] == {
        "section": 1,
        "line_start": 1,
        "line_end": 3,
    }
    assert markdown_blocks[1]["locator"] == {
        "section": 2,
        "line_start": 5,
        "line_end": 6,
    }
    assert plain_blocks[0]["locator"] == {
        "paragraph": 0,
        "line_start": 1,
        "line_end": 2,
    }
    assert plain_blocks[1]["locator"] == {
        "paragraph": 1,
        "line_start": 4,
        "line_end": 4,
    }

def test_markdown_import_preserves_deep_heading_hierarchy_and_standalone_anchors(tmp_path) -> None:
    markdown = """# 合订版

<a id="part-one"></a>

## 第一部分

### 第一册

#### 第一章

##### 规则 A
正文
"""

    with TestClient(create_app(data_dir=tmp_path)) as client:
        item = client.post(
            "/api/documents/import",
            files={"file": ("collection.md", markdown.encode("utf-8"), "text/markdown")},
        ).json()["data"]
        content = client.get(f"/api/documents/{item['id']}/content").json()["data"]

    blocks = content["blocks"]
    assert [block["data"]["heading_level"] for block in blocks] == [1, 2, 3, 4, 5]
    assert blocks[1]["data"]["anchor"] == "part-one"
    assert blocks[1]["data"]["parent_block_key"] == blocks[0]["block_key"]
    assert blocks[4]["data"]["heading_path"] == [
        "合订版",
        "第一部分",
        "第一册",
        "第一章",
        "规则 A",
    ]
    assert content["document"]["metadata"]["max_heading_depth"] == 5
    assert '<a id="part-one"></a>' not in blocks[0]["text"]


def test_imports_jupyter_notebook_with_markdown_code_and_outputs(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        response = client.post(
            "/api/documents/import",
            files={"file": ("experiment.ipynb", make_notebook(), "application/x-ipynb+json")},
        )

        assert response.status_code == 201, response.text
        item = response.json()["data"]
        assert item["format"] == "ipynb"
        content = client.get(f"/api/documents/{item['id']}/content").json()["data"]

    assert [block["block_type"] for block in content["blocks"]] == [
        "notebook_markdown",
        "notebook_code",
    ]
    assert content["blocks"][1]["data"]["execution_count"] == 7
    assert content["blocks"][1]["data"]["outputs"][0]["text"] == "0.95"
    assert "score = 0.95" in item["body"]
    assert "0.95" in item["body"]


def test_serves_the_original_pdf_inline_and_accepts_pages_without_text(tmp_path) -> None:
    original = make_blank_pdf()
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("scan.pdf", original, "application/pdf")},
        )
        assert imported.status_code == 201, imported.text
        item = imported.json()["data"]
        raw = client.get(f"/api/documents/{item['id']}/file")

    assert raw.status_code == 200
    assert raw.content == original
    assert raw.headers["content-type"].startswith("application/pdf")
    assert raw.headers["content-disposition"].startswith("inline")


def test_xlsx_parser_preserves_merged_cells_and_visual_formatting(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={
                "file": (
                    "styled.xlsx",
                    make_styled_xlsx(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )
        assert imported.status_code == 201, imported.text
        item = imported.json()["data"]
        block = client.get(f"/api/documents/{item['id']}/content").json()["data"]["blocks"][0]

    assert block["data"]["merged_ranges"] == ["A1:B1"]
    assert block["data"]["column_widths"]["A"] == 24.0
    assert block["data"]["row_heights"]["1"] == 30.0
    cell = next(item for item in block["data"]["cells"] if item["address"] == "A1")
    assert cell["style"]["font"]["bold"] is True
    assert cell["style"]["font"]["color"] == "FFFFFF"
    assert cell["style"]["fill"] == "2F5597"
    assert cell["style"]["alignment"]["horizontal"] == "center"


def test_pptx_parser_preserves_shape_position_and_text_formatting(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={
                "file": (
                    "styled.pptx",
                    make_styled_pptx(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert imported.status_code == 201, imported.text
        item = imported.json()["data"]
        block = client.get(f"/api/documents/{item['id']}/content").json()["data"]["blocks"][0]

    element = block["data"]["elements"][0]
    assert 0.09 < element["layout"]["left"] < 0.11
    assert 0.19 < element["layout"]["top"] < 0.21
    assert element["style"]["font_size"] == 28.0
    assert element["style"]["bold"] is True
    assert element["style"]["color"] == "2F5597"


def test_full_text_search_and_highlight_use_document_body(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("retrieval.md", "# 检索\nBM25 是稀疏检索基线".encode(), "text/markdown")},
        ).json()["data"]
        search = client.get("/api/search", params={"q": "BM25"})
        highlight = client.post(
            f"/api/documents/{imported['id']}/highlights",
            json={"quote": "BM25 是稀疏检索基线", "note": "加入知识点"},
        )

    assert search.status_code == 200
    assert search.json()["data"][0]["document_id"] == imported["id"]
    assert "<mark>BM25</mark>" in search.json()["data"][0]["snippet"]
    assert highlight.status_code == 201


def test_duplicate_document_is_deduplicated(tmp_path) -> None:
    content = "same content".encode()
    with TestClient(create_app(data_dir=tmp_path)) as client:
        first = client.post("/api/documents/import", files={"file": ("a.txt", content, "text/plain")})
        second = client.post("/api/documents/import", files={"file": ("b.txt", content, "text/plain")})

    assert first.status_code == 201
    assert second.status_code == 200
    assert second.json()["data"]["id"] == first.json()["data"]["id"]
    assert second.json()["meta"]["deduplicated"] is True


def test_reimporting_a_deleted_document_restores_it_to_the_shelf(tmp_path) -> None:
    content = "restore the same source".encode()
    with TestClient(create_app(data_dir=tmp_path)) as client:
        first = client.post("/api/documents/import", files={"file": ("original.txt", content, "text/plain")})
        document_id = first.json()["data"]["id"]
        deleted = client.delete(f"/api/documents/{document_id}")
        second = client.post("/api/documents/import", files={"file": ("original.txt", content, "text/plain")})
        visible = client.get("/api/documents").json()["data"]

    assert deleted.status_code == 204
    assert second.status_code == 200
    assert second.json()["meta"]["deduplicated"] is True
    assert second.json()["data"]["id"] == document_id
    assert second.json()["data"]["deleted_at"] is None
    assert [item["id"] for item in visible] == [document_id]


def test_rejects_unsupported_file_without_copying_it(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        response = client.post(
            "/api/documents/import", files={"file": ("payload.exe", b"MZ", "application/octet-stream")}
        )

    assert response.status_code == 415
    assert response.json()["error"]["code"] == "UNSUPPORTED_DOCUMENT"
    assert not list((tmp_path / "documents").glob("*"))


def test_document_can_be_renamed_revised_and_searched(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("draft.txt", b"Original retrieval note", "text/plain")},
        ).json()["data"]
        content = client.get(f"/api/documents/{imported['id']}/content").json()["data"]
        block = content["blocks"][0]
        renamed = client.patch(
            f"/api/documents/{imported['id']}", json={"title": "检索修订稿"}
        )
        revised = client.post(
            f"/api/documents/{imported['id']}/revisions",
            json={
                "block_key": block["block_key"],
                "before": {"text": block["text"]},
                "after": {"text": "Hybrid retrieval combines sparse and dense signals"},
            },
        )
        search = client.get("/api/search", params={"q": "Hybrid retrieval"})

    assert renamed.status_code == 200
    assert renamed.json()["data"]["title"] == "检索修订稿"
    assert revised.status_code == 201
    assert revised.json()["data"]["block"]["text"].startswith("Hybrid retrieval")
    assert search.json()["data"][0]["document_id"] == imported["id"]


def test_document_favorite_and_pin_are_persisted_without_renaming(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("favorite.txt", b"Keep this source", "text/plain")},
        ).json()["data"]
        updated = client.patch(
            f"/api/documents/{imported['id']}",
            json={"favorite": True, "pinned": True},
        )
        restored = client.get(f"/api/documents/{imported['id']}")

    assert updated.status_code == 200
    assert updated.json()["data"]["title"] == imported["title"]
    assert restored.json()["data"]["metadata"]["favorite"] is True
    assert restored.json()["data"]["metadata"]["pinned"] is True


def test_document_annotations_support_create_update_and_delete(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("notes.md", b"# Evidence\nQuoted claim", "text/markdown")},
        ).json()["data"]
        block = client.get(f"/api/documents/{imported['id']}/content").json()["data"]["blocks"][0]
        created = client.post(
            f"/api/documents/{imported['id']}/annotations",
            json={
                "block_key": block["block_key"],
                "kind": "ellipse",
                "locator": block["locator"],
                "quote": "Quoted claim",
                "note": "需要核验",
                "color": "coral",
                "geometry": {"x": 0.12, "y": 0.2, "width": 0.4, "height": 0.2},
            },
        )
        annotation_id = created.json()["data"]["id"]
        updated = client.patch(
            f"/api/documents/{imported['id']}/annotations/{annotation_id}",
            json={"note": "已核验"},
        )
        listed = client.get(f"/api/documents/{imported['id']}/annotations")
        deleted = client.delete(
            f"/api/documents/{imported['id']}/annotations/{annotation_id}"
        )
        empty = client.get(f"/api/documents/{imported['id']}/annotations")

    assert created.status_code == 201
    assert updated.json()["data"]["note"] == "已核验"
    assert listed.json()["data"][0]["geometry"]["width"] == 0.4
    assert deleted.status_code == 204
    assert empty.json()["data"] == []


def test_document_revisions_support_undo_redo_and_clear_the_redo_branch(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("history.txt", b"Original content", "text/plain")},
        ).json()["data"]
        block = client.get(
            f"/api/documents/{imported['id']}/content"
        ).json()["data"]["blocks"][0]
        revised = client.post(
            f"/api/documents/{imported['id']}/revisions",
            json={
                "block_key": block["block_key"],
                "before": {"text": "Original content"},
                "after": {"text": "Revised content"},
            },
        )
        state_after_edit = client.get(
            f"/api/documents/{imported['id']}/revisions"
        )
        undone = client.post(f"/api/documents/{imported['id']}/revisions/undo")
        state_after_undo = client.get(
            f"/api/documents/{imported['id']}/revisions"
        )
        redone = client.post(f"/api/documents/{imported['id']}/revisions/redo")

        client.post(f"/api/documents/{imported['id']}/revisions/undo")
        branched = client.post(
            f"/api/documents/{imported['id']}/revisions",
            json={
                "block_key": block["block_key"],
                "before": {"text": "Original content"},
                "after": {"text": "Branched content"},
            },
        )
        state_after_branch = client.get(
            f"/api/documents/{imported['id']}/revisions"
        )

    assert revised.status_code == 201
    assert state_after_edit.status_code == 200
    assert state_after_edit.json()["data"] == {"can_undo": True, "can_redo": False}
    assert undone.status_code == 200
    assert undone.json()["data"]["block"]["text"] == "Original content"
    assert state_after_undo.json()["data"] == {"can_undo": False, "can_redo": True}
    assert redone.status_code == 200
    assert redone.json()["data"]["block"]["text"] == "Revised content"
    assert branched.status_code == 201
    assert state_after_branch.json()["data"] == {"can_undo": True, "can_redo": False}


def test_document_soft_delete_restore_and_course_isolation(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("private.txt", b"Private course source", "text/plain")},
        ).json()["data"]
        stored = tmp_path / imported["stored_path"]
        deleted = client.delete(f"/api/documents/{imported['id']}")
        visible = client.get("/api/documents").json()["data"]
        trash = client.get("/api/documents", params={"include_deleted": "true"}).json()["data"]
        other = client.post("/api/courses", json={"title": "Other"}).json()["data"]
        client.post(f"/api/courses/{other['id']}/activate")
        isolated = client.get(f"/api/documents/{imported['id']}")
        client.post("/api/courses/1/activate")
        restored = client.post(f"/api/documents/{imported['id']}/restore")

    assert deleted.status_code == 204
    assert visible == []
    assert trash[0]["id"] == imported["id"]
    assert trash[0]["deleted_at"] is not None
    assert stored.is_file()
    assert isolated.status_code == 404
    assert restored.status_code == 200
    assert restored.json()["data"]["deleted_at"] is None


def test_exports_saved_markdown_and_word_content_with_explicit_pdf_fallback(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        markdown = client.post(
            "/api/documents/import",
            files={"file": ("notes.md", "# 旧标题\n旧内容".encode("utf-8"), "text/markdown")},
        ).json()["data"]
        block = client.get(f"/api/documents/{markdown['id']}/content").json()["data"]["blocks"][0]
        client.post(
            f"/api/documents/{markdown['id']}/revisions",
            json={
                "block_key": block["block_key"], "before": {"text": block["text"]},
                "after": {"text": "# 新标题\n已保存的内容"},
            },
        )
        markdown_source = client.post(f"/api/documents/{markdown['id']}/export", json={"format": "source"})
        markdown_pdf = client.post(f"/api/documents/{markdown['id']}/export", json={"format": "pdf"})
        word = client.post(
            "/api/documents/import",
            files={"file": ("lesson.docx", make_docx("Word saved content"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
        ).json()["data"]
        word_source = client.post(f"/api/documents/{word['id']}/export", json={"format": "source"})

    assert markdown_source.status_code == 200
    assert "已保存的内容" in markdown_source.content.decode("utf-8")
    assert "filename*=UTF-8''" in markdown_source.headers["content-disposition"]
    assert markdown_pdf.status_code == 200
    assert markdown_pdf.content.startswith(b"%PDF")
    assert word_source.status_code == 200
    exported_word = Document(BytesIO(word_source.content))
    assert any("Word saved content" in paragraph.text for paragraph in exported_word.paragraphs)


def test_exports_saved_csv_cells_as_a_real_csv_source(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        imported = client.post(
            "/api/documents/import",
            files={"file": ("metrics.csv", b"name,score\nRRF,0.95\n", "text/csv")},
        ).json()["data"]
        block = client.get(f"/api/documents/{imported['id']}/content").json()["data"]["blocks"][0]
        changed_cells = [
            {**cell, "value": "Hybrid RRF" if cell["address"] == "A2" else cell["value"]}
            for cell in block["data"]["cells"]
        ]
        revised = client.post(
            f"/api/documents/{imported['id']}/revisions",
            json={
                "block_key": block["block_key"],
                "before": {"text": block["text"], "data": block["data"]},
                "after": {
                    "text": "A1: name\nB1: score\nA2: Hybrid RRF\nB2: 0.95",
                    "data": {**block["data"], "cells": changed_cells},
                },
            },
        )
        exported = client.post(f"/api/documents/{imported['id']}/export", json={"format": "source"})

    assert revised.status_code == 201, revised.text
    assert exported.status_code == 200, exported.text
    assert exported.headers["content-type"].startswith("text/csv")
    rows = list(csv.reader(StringIO(exported.content.decode("utf-8-sig"))))
    assert rows == [["name", "score"], ["Hybrid RRF", "0.95"]]


def test_exports_editable_text_and_tsv_using_their_original_formats(tmp_path) -> None:
    with TestClient(create_app(data_dir=tmp_path)) as client:
        python_item = client.post(
            "/api/documents/import",
            files={"file": ("analysis.py", b"print('before')\n", "text/x-python")},
        ).json()["data"]
        python_block = client.get(f"/api/documents/{python_item['id']}/content").json()["data"]["blocks"][0]
        client.post(
            f"/api/documents/{python_item['id']}/revisions",
            json={
                "block_key": python_block["block_key"],
                "before": {"text": python_block["text"]},
                "after": {"text": "print('saved')"},
            },
        )
        python_export = client.post(f"/api/documents/{python_item['id']}/export", json={"format": "source"})

        tsv_item = client.post(
            "/api/documents/import",
            files={"file": ("metrics.tsv", b"name\tscore\nRRF\t0.95\n", "text/tab-separated-values")},
        ).json()["data"]
        tsv_export = client.post(f"/api/documents/{tsv_item['id']}/export", json={"format": "source"})

    assert python_export.status_code == 200
    assert "analysis.py" in python_export.headers["content-disposition"]
    assert python_export.content.decode("utf-8").strip() == "print('saved')"
    assert tsv_export.status_code == 200
    assert "metrics.tsv" in tsv_export.headers["content-disposition"]
    assert tsv_export.headers["content-type"].startswith("text/tab-separated-values")
    assert list(csv.reader(StringIO(tsv_export.content.decode("utf-8-sig")), delimiter="\t")) == [
        ["name", "score"],
        ["RRF", "0.95"],
    ]
