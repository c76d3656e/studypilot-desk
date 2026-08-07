from __future__ import annotations

import csv
import hashlib
import json
import re
from dataclasses import dataclass
from datetime import date, datetime, time
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas

from ..db import Database
from ..errors import AppError
from ..repository import as_dict


# Compatibility marker for older callers. Document imports no longer have a fixed byte ceiling.
MAX_DOCUMENT_BYTES: int | None = None
TEXT_SUFFIXES = {
    ".txt", ".json", ".yaml", ".yml", ".xml", ".html", ".css", ".js", ".jsx",
    ".ts", ".tsx", ".py", ".java", ".c", ".cpp", ".h", ".hpp", ".sql", ".log",
    ".ini", ".toml",
}
SUPPORTED = {".pdf", ".docx", ".md", ".markdown", ".csv", ".tsv", ".xlsx", ".pptx", ".ipynb", *TEXT_SUFFIXES}
FORMAT_BY_SUFFIX = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".md": "markdown",
    ".markdown": "markdown",
    ".csv": "csv",
    ".tsv": "csv",
    **{suffix: "text" for suffix in TEXT_SUFFIXES},
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".ipynb": "ipynb",
}


@dataclass(frozen=True)
class ParsedBlock:
    block_key: str
    block_type: str
    ordinal: int
    locator: dict[str, Any]
    text: str
    data: dict[str, Any]


@dataclass(frozen=True)
class ParsedDocument:
    format: str
    structure: dict[str, Any]
    metadata: dict[str, Any]
    blocks: list[ParsedBlock]

    @property
    def body(self) -> str:
        return "\n\n".join(block.text.strip() for block in self.blocks if block.text.strip())


class DocumentService:
    def __init__(self, database: Database, data_dir: Path) -> None:
        self.database = database
        self.data_dir = data_dir
        self.documents_dir = data_dir / "documents"

    def _course_id(self) -> int:
        with self.database.connect() as connection:
            row = connection.execute(
                "SELECT value_json FROM settings WHERE key = 'active_course'"
            ).fetchone()
        return int(json.loads(row[0]))

    def parse(self, filename: str, content: bytes) -> ParsedDocument:
        suffix = Path(filename).suffix.lower()
        if suffix not in SUPPORTED:
            raise AppError(
                "UNSUPPORTED_DOCUMENT",
                "仅支持 PDF、DOCX、Markdown、CSV、常见文本/代码、XLSX、PPTX 和 Jupyter Notebook",
                415,
            )
        if not content:
            raise AppError("EMPTY_DOCUMENT", "文件内容为空", 422)

        parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".md": self._parse_markdown,
            ".markdown": self._parse_markdown,
            ".csv": self._parse_csv,
            ".tsv": self._parse_csv,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".ipynb": self._parse_ipynb,
            **{text_suffix: self._parse_text for text_suffix in TEXT_SUFFIXES},
        }
        try:
            parsed = parsers[suffix](content)
        except AppError:
            raise
        except Exception as exc:
            if suffix == ".pdf":
                return ParsedDocument(
                    format="pdf",
                    structure={"pages": []},
                    metadata={"pages": 0, "parse_fallback": True},
                    blocks=[ParsedBlock(
                        block_key="original:1", block_type="original", ordinal=0, locator={},
                        text="PDF 无法结构化解析，已保留原文件，可直接阅读。", data={"parse_fallback": True},
                    )],
                )
            raise AppError(
                "DOCUMENT_PARSE_FAILED",
                f"{FORMAT_BY_SUFFIX[suffix].upper()} 文档解析失败",
                422,
            ) from exc

        return parsed

    def _parse_pdf(self, content: bytes) -> ParsedDocument:
        reader = PdfReader(BytesIO(content))
        blocks = [
            ParsedBlock(
                block_key=f"page:{index}",
                block_type="page",
                ordinal=index - 1,
                locator={"page": index},
                text=(page.extract_text() or "").strip(),
                data={"page": index},
            )
            for index, page in enumerate(reader.pages, start=1)
        ]
        return ParsedDocument(
            format="pdf",
            structure={"pages": [{"page": index} for index in range(1, len(reader.pages) + 1)]},
            metadata={"pages": len(reader.pages)},
            blocks=blocks,
        )

    def _parse_docx(self, content: bytes) -> ParsedDocument:
        document = Document(BytesIO(content))
        blocks: list[ParsedBlock] = []
        outline: list[dict[str, Any]] = []
        for index, paragraph in enumerate(document.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style.name if paragraph.style else "Normal"
            block = ParsedBlock(
                block_key=f"paragraph:{index}",
                block_type="paragraph",
                ordinal=len(blocks),
                locator={"paragraph": index},
                text=text,
                data={"style": style},
            )
            blocks.append(block)
            if style.lower().startswith("heading"):
                outline.append({"block_key": block.block_key, "title": text, "style": style})
        for table_index, table in enumerate(document.tables):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            text = "\n".join(" | ".join(cells) for cells in rows if any(cells)).strip()
            if not text:
                continue
            blocks.append(
                ParsedBlock(
                    block_key=f"table:{table_index}",
                    block_type="table",
                    ordinal=len(blocks),
                    locator={"table": table_index},
                    text=text,
                    data={"rows": rows},
                )
            )
        return ParsedDocument(
            format="docx",
            structure={"outline": outline, "paragraphs": len(document.paragraphs), "tables": len(document.tables)},
            metadata={"paragraphs": len(document.paragraphs), "tables": len(document.tables)},
            blocks=blocks,
        )

    def _parse_markdown(self, content: bytes) -> ParsedDocument:
        text = self._decode_text(content)
        source_lines = text.splitlines()
        heading_anchors: dict[int, str] = {}
        discarded_anchor_lines: set[int] = set()
        standalone_anchor = re.compile(
            r"""^\s*<a\s+(?:id|name)=["']([^"']+)["']\s*>\s*</a>\s*$""",
            flags=re.IGNORECASE,
        )
        for source_index, line in enumerate(source_lines):
            anchor_match = standalone_anchor.match(line)
            if not anchor_match:
                continue
            candidate = source_index + 1
            while candidate < len(source_lines):
                stripped = source_lines[candidate].strip()
                if not stripped or (stripped.startswith("<!--") and stripped.endswith("-->")):
                    candidate += 1
                    continue
                if re.match(r"^\s{0,3}#{1,6}\s+", source_lines[candidate]):
                    heading_anchors[candidate + 1] = anchor_match.group(1)
                    discarded_anchor_lines.add(source_index + 1)
                break

        sections: list[tuple[str, int, str, list[tuple[int, str]]]] = []
        title = "正文"
        heading_level = 0
        anchor = ""
        lines: list[tuple[int, str]] = []
        for line_number, line in enumerate(source_lines, start=1):
            if line_number in discarded_anchor_lines:
                continue
            heading = re.match(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$", line)
            if heading:
                if lines or not sections:
                    sections.append((title, heading_level, anchor, lines))
                title = heading.group(2).strip()
                heading_level = len(heading.group(1))
                anchor = heading_anchors.get(line_number, "")
                lines = [(line_number, line)]
            else:
                lines.append((line_number, line))
        if lines:
            sections.append((title, heading_level, anchor, lines))

        blocks: list[ParsedBlock] = []
        heading_stack: dict[int, ParsedBlock] = {}
        for section_index, (section_title, level, section_anchor, section_lines) in enumerate(sections):
            meaningful = [item for item in section_lines if item[1].strip()]
            if not meaningful:
                continue
            line_start = meaningful[0][0]
            line_end = meaningful[-1][0]
            rendered = "\n".join(
                line for number, line in section_lines if line_start <= number <= line_end
            ).strip()
            parent = next(
                (heading_stack[parent_level] for parent_level in range(level - 1, 0, -1) if parent_level in heading_stack),
                None,
            ) if level else None
            heading_path = [
                heading_stack[parent_level].data["title"]
                for parent_level in sorted(heading_stack)
                if parent_level < level
            ] + ([section_title] if level else [])
            block = ParsedBlock(
                block_key=f"section:{section_index}",
                block_type="section",
                ordinal=len(blocks),
                locator={
                    "section": section_index,
                    "line_start": line_start,
                    "line_end": line_end,
                },
                text=rendered,
                data={
                    "title": section_title,
                    "heading_level": level,
                    "parent_block_key": parent.block_key if parent else "",
                    "heading_path": heading_path,
                    "anchor": section_anchor,
                },
            )
            blocks.append(block)
            if level:
                heading_stack = {
                    stack_level: stack_block
                    for stack_level, stack_block in heading_stack.items()
                    if stack_level < level
                }
                heading_stack[level] = block
        max_heading_depth = max((int(block.data.get("heading_level") or 0) for block in blocks), default=0)
        return ParsedDocument(
            format="markdown",
            structure={
                "sections": [
                    {
                        "block_key": block.block_key,
                        "title": block.data["title"],
                        "heading_level": block.data["heading_level"],
                        "parent_block_key": block.data["parent_block_key"],
                        "anchor": block.data["anchor"],
                    }
                    for block in blocks
                ]
            },
            metadata={
                "characters": len(text),
                "sections": len(blocks),
                "max_heading_depth": max_heading_depth,
                "large_document": len(text) >= 200_000 or len(blocks) >= 120,
            },
            blocks=blocks,
        )
    def _parse_text(self, content: bytes) -> ParsedDocument:
        text = self._decode_text(content)
        paragraphs: list[list[tuple[int, str]]] = []
        current: list[tuple[int, str]] = []
        for line_number, line in enumerate(text.splitlines(), start=1):
            if line.strip():
                current.append((line_number, line))
            elif current:
                paragraphs.append(current)
                current = []
        if current:
            paragraphs.append(current)

        blocks = [
            ParsedBlock(
                block_key=f"paragraph:{index}",
                block_type="paragraph",
                ordinal=index,
                locator={
                    "paragraph": index,
                    "line_start": paragraph[0][0],
                    "line_end": paragraph[-1][0],
                },
                text="\n".join(line for _, line in paragraph).strip(),
                data={},
            )
            for index, paragraph in enumerate(paragraphs)
        ]
        return ParsedDocument(
            format="text",
            structure={"paragraphs": len(blocks)},
            metadata={"characters": len(text), "paragraphs": len(blocks)},
            blocks=blocks,
        )

    def _parse_csv(self, content: bytes) -> ParsedDocument:
        text = self._decode_text(content)
        sample = text[:16384]
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",\t;")
            rows = list(csv.reader(StringIO(text, newline=""), dialect))
            delimiter = dialect.delimiter
        except csv.Error:
            delimiter = "\t" if "\t" in sample else ","
            rows = list(csv.reader(StringIO(text, newline=""), delimiter=delimiter))
        while rows and not any(value for value in rows[-1]):
            rows.pop()
        max_columns = max((len(row) for row in rows), default=1)
        max_rows = max(1, len(rows))

        def column_name(column: int) -> str:
            result = ""
            while column > 0:
                column, remainder = divmod(column - 1, 26)
                result = chr(65 + remainder) + result
            return result

        cells = [
            {
                "address": f"{column_name(column_index)}{row_index}",
                "row": row_index,
                "column": column_index,
                "value": row[column_index - 1] if column_index <= len(row) else "",
                "formula": False,
                "number_format": "General",
                "style": {},
            }
            for row_index, row in enumerate(rows, start=1)
            for column_index in range(1, max_columns + 1)
        ]
        dimensions = f"A1:{column_name(max_columns)}{max_rows}"
        block = ParsedBlock(
            block_key="sheet:0",
            block_type="sheet",
            ordinal=0,
            locator={"sheet": "CSV", "range": dimensions},
            text="\n".join(delimiter.join(row) for row in rows),
            data={
                "title": "CSV",
                "cells": cells,
                "dimensions": dimensions,
                "delimiter": delimiter,
                "merged_ranges": [],
                "column_widths": {},
                "row_heights": {},
            },
        )
        return ParsedDocument(
            format="csv",
            structure={"sheets": [{"block_key": block.block_key, "title": "CSV", "dimensions": dimensions}]},
            metadata={"rows": len(rows), "columns": max_columns, "delimiter": delimiter},
            blocks=[block],
        )

    def _parse_xlsx(self, content: bytes) -> ParsedDocument:
        workbook = load_workbook(BytesIO(content), read_only=False, data_only=False)
        blocks: list[ParsedBlock] = []
        sheets: list[dict[str, Any]] = []
        try:
            for sheet_index, worksheet in enumerate(workbook.worksheets):
                cells: list[dict[str, Any]] = []
                text_parts: list[str] = []
                for row in worksheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        value = self._json_value(cell.value)
                        cells.append(
                            {
                                "address": cell.coordinate,
                                "row": cell.row,
                                "column": cell.column,
                                "value": value,
                                "formula": isinstance(cell.value, str) and cell.value.startswith("="),
                                "number_format": cell.number_format,
                                "style": self._xlsx_style(cell),
                            }
                        )
                        text_parts.append(f"{cell.coordinate}: {value}")
                if not cells:
                    continue
                dimensions = worksheet.calculate_dimension()
                block = ParsedBlock(
                    block_key=f"sheet:{sheet_index}",
                    block_type="sheet",
                    ordinal=len(blocks),
                    locator={"sheet": worksheet.title, "range": dimensions},
                    text="\n".join(text_parts),
                    data={
                        "title": worksheet.title,
                        "cells": cells,
                        "dimensions": dimensions,
                        "merged_ranges": [str(item) for item in worksheet.merged_cells.ranges],
                        "column_widths": {
                            str(name): float(dimension.width)
                            for name, dimension in worksheet.column_dimensions.items()
                            if dimension.width is not None
                        },
                        "row_heights": {
                            str(index): float(dimension.height)
                            for index, dimension in worksheet.row_dimensions.items()
                            if dimension.height is not None
                        },
                    },
                )
                blocks.append(block)
                sheets.append({"block_key": block.block_key, "title": worksheet.title, "dimensions": dimensions})
        finally:
            workbook.close()
        return ParsedDocument(
            format="xlsx",
            structure={"sheets": sheets},
            metadata={"sheets": len(workbook.sheetnames)},
            blocks=blocks,
        )

    def _parse_ipynb(self, content: bytes) -> ParsedDocument:
        notebook = json.loads(self._decode_text(content))
        if not isinstance(notebook, dict) or not isinstance(notebook.get("cells"), list):
            raise AppError("DOCUMENT_PARSE_FAILED", "Jupyter Notebook 结构无效", 422)
        blocks: list[ParsedBlock] = []
        outline: list[dict[str, Any]] = []
        for index, cell in enumerate(notebook["cells"]):
            if not isinstance(cell, dict):
                continue
            cell_type = str(cell.get("cell_type") or "raw")
            source = self._notebook_text(cell.get("source"))
            outputs = [
                output
                for item in (cell.get("outputs") or [])
                if isinstance(item, dict)
                for output in [self._notebook_output(item)]
                if output["text"] or output.get("html")
            ]
            searchable = "\n\n".join(
                part for part in [source, *(str(item["text"]) for item in outputs)] if part
            )
            block = ParsedBlock(
                block_key=f"cell:{index}",
                block_type=f"notebook_{cell_type}",
                ordinal=len(blocks),
                locator={"cell": index + 1, "cell_type": cell_type},
                text=searchable,
                data={
                    "cell_type": cell_type,
                    "source": source,
                    "execution_count": cell.get("execution_count"),
                    "outputs": outputs,
                    "metadata": cell.get("metadata") if isinstance(cell.get("metadata"), dict) else {},
                },
            )
            blocks.append(block)
            if cell_type == "markdown":
                heading = re.search(r"^\s*#{1,6}\s+(.+)$", source, flags=re.MULTILINE)
                if heading:
                    outline.append({"block_key": block.block_key, "title": heading.group(1).strip()})
        return ParsedDocument(
            format="ipynb",
            structure={
                "cells": [
                    {"block_key": block.block_key, "cell": block.locator["cell"], "cell_type": block.data["cell_type"]}
                    for block in blocks
                ],
                "outline": outline,
            },
            metadata={
                "cells": len(blocks),
                "nbformat": notebook.get("nbformat"),
                "kernel": ((notebook.get("metadata") or {}).get("kernelspec") or {}).get("display_name", ""),
            },
            blocks=blocks,
        )

    @staticmethod
    def _notebook_text(value: Any) -> str:
        if isinstance(value, list):
            return "".join(str(item) for item in value)
        return str(value or "")

    @classmethod
    def _notebook_output(cls, output: dict[str, Any]) -> dict[str, Any]:
        output_type = str(output.get("output_type") or "output")
        data = output.get("data") if isinstance(output.get("data"), dict) else {}
        text = cls._notebook_text(
            output.get("text")
            or data.get("text/plain")
            or output.get("traceback")
            or ""
        )
        if output_type == "error" and not text:
            text = f"{output.get('ename', 'Error')}: {output.get('evalue', '')}".strip()
        return {
            "output_type": output_type,
            "text": text,
            "html": cls._notebook_text(data.get("text/html")),
        }

    @staticmethod
    def _xlsx_color(color: Any) -> str:
        if color is None:
            return ""
        if getattr(color, "type", "") == "rgb" and color.rgb:
            value = str(color.rgb)
            return value[-6:].upper()
        return ""

    @classmethod
    def _xlsx_style(cls, cell: Any) -> dict[str, Any]:
        return {
            "font": {
                "name": cell.font.name or "",
                "size": float(cell.font.sz) if cell.font.sz is not None else None,
                "bold": bool(cell.font.bold),
                "italic": bool(cell.font.italic),
                "color": cls._xlsx_color(cell.font.color),
            },
            "fill": cls._xlsx_color(cell.fill.fgColor) if cell.fill.fill_type else "",
            "alignment": {
                "horizontal": cell.alignment.horizontal or "",
                "vertical": cell.alignment.vertical or "",
                "wrap_text": bool(cell.alignment.wrap_text),
            },
        }

    def _parse_pptx(self, content: bytes) -> ParsedDocument:
        presentation = Presentation(BytesIO(content))
        slide_width = float(presentation.slide_width or 1)
        slide_height = float(presentation.slide_height or 1)
        blocks: list[ParsedBlock] = []
        slides: list[dict[str, Any]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            elements: list[dict[str, Any]] = []
            text_parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if not text:
                    continue
                paragraph = shape.text_frame.paragraphs[0] if getattr(shape, "has_text_frame", False) and shape.text_frame.paragraphs else None
                run = paragraph.runs[0] if paragraph is not None and paragraph.runs else None
                color = ""
                if run is not None:
                    try:
                        color = str(run.font.color.rgb or "")
                    except (AttributeError, TypeError):
                        color = ""
                elements.append(
                    {
                        "shape_id": shape.shape_id,
                        "name": shape.name,
                        "text": text,
                        "layout": {
                            "left": float(shape.left or 0) / slide_width,
                            "top": float(shape.top or 0) / slide_height,
                            "width": float(shape.width or 0) / slide_width,
                            "height": float(shape.height or 0) / slide_height,
                            "rotation": float(shape.rotation or 0),
                        },
                        "style": {
                            "font_size": float(run.font.size.pt) if run is not None and run.font.size is not None else None,
                            "bold": bool(run.font.bold) if run is not None else False,
                            "italic": bool(run.font.italic) if run is not None else False,
                            "color": color,
                            "align": paragraph.alignment.name.lower() if paragraph is not None and paragraph.alignment is not None else "left",
                        },
                    }
                )
                text_parts.append(text)
            if not text_parts:
                continue
            block = ParsedBlock(
                block_key=f"slide:{slide_index}",
                block_type="slide",
                ordinal=len(blocks),
                locator={"slide": slide_index},
                text="\n".join(text_parts),
                data={
                    "elements": elements,
                    "slide_size": {"width": slide_width, "height": slide_height},
                },
            )
            blocks.append(block)
            slides.append({"block_key": block.block_key, "slide": slide_index, "title": text_parts[0]})
        return ParsedDocument(
            format="pptx",
            structure={"slides": slides},
            metadata={"slides": len(presentation.slides)},
            blocks=blocks,
        )

    @staticmethod
    def _json_value(value: Any) -> str | int | float | bool | None:
        if value is None or isinstance(value, (str, int, float, bool)):
            return value
        if isinstance(value, (date, datetime, time)):
            return value.isoformat()
        return str(value)

    @staticmethod
    def _decode_text(content: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-8", "gb18030"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise AppError("DOCUMENT_ENCODING", "文本编码无法识别", 422)

    def import_bytes(
        self,
        filename: str,
        media_type: str,
        content: bytes,
        source_created_at: str | None = None,
    ) -> tuple[dict, bool]:
        parsed = self.parse(filename, content)
        course_id = self._course_id()
        digest = hashlib.sha256(content).hexdigest()
        source_created_at = self._normalize_source_created_at(source_created_at)
        with self.database.connect() as connection:
            existing = connection.execute(
                "SELECT * FROM documents WHERE course_id = ? AND sha256 = ?",
                (course_id, digest),
            ).fetchone()
        if existing:
            if existing["deleted_at"] or (source_created_at and not existing["source_created_at"]):
                with self.database.connect() as connection:
                    connection.execute(
                        """UPDATE documents
                        SET deleted_at = NULL,
                            source_created_at = COALESCE(source_created_at, ?),
                            updated_at = CURRENT_TIMESTAMP
                        WHERE id = ?""",
                        (source_created_at, existing["id"]),
                    )
                    existing = connection.execute(
                        "SELECT * FROM documents WHERE id = ?", (existing["id"],)
                    ).fetchone()
            return self._serialize(existing), True

        self.documents_dir.mkdir(parents=True, exist_ok=True)
        suffix = Path(filename).suffix.lower()
        stored_name = f"{digest[:16]}{suffix}"
        stored = self.documents_dir / stored_name
        temporary = stored.with_suffix(stored.suffix + ".tmp")
        temporary.write_bytes(content)
        temporary.replace(stored)
        relative_path = stored.relative_to(self.data_dir).as_posix()

        try:
            with self.database.connect() as connection:
                cursor = connection.execute(
                    """INSERT INTO documents(
                        course_id, title, filename, stored_path, media_type, sha256,
                        body, metadata_json, format, status, structure_json,
                        source_created_at, updated_at
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, 'ready', ?, ?, CURRENT_TIMESTAMP)""",
                    (
                        course_id,
                        Path(filename).stem,
                        filename,
                        relative_path,
                        media_type or "application/octet-stream",
                        digest,
                        parsed.body,
                        json.dumps(parsed.metadata, ensure_ascii=False),
                        parsed.format,
                        json.dumps(parsed.structure, ensure_ascii=False),
                        source_created_at,
                    ),
                )
                document_id = int(cursor.lastrowid)
                connection.executemany(
                    """INSERT INTO document_blocks(
                        document_id, block_key, block_type, ordinal, locator_json, text, data_json
                    ) VALUES (?, ?, ?, ?, ?, ?, ?)""",
                    [
                        (
                            document_id,
                            block.block_key,
                            block.block_type,
                            block.ordinal,
                            json.dumps(block.locator, ensure_ascii=False),
                            block.text,
                            json.dumps(block.data, ensure_ascii=False),
                        )
                        for block in parsed.blocks
                    ],
                )
                row = connection.execute(
                    "SELECT * FROM documents WHERE id = ?", (document_id,)
                ).fetchone()
        except Exception:
            stored.unlink(missing_ok=True)
            raise
        return self._serialize(row), False

    @staticmethod
    def _normalize_source_created_at(value: str | None) -> str | None:
        timestamp = (value or "").strip()
        if not timestamp:
            return None
        if len(timestamp) > 64:
            raise AppError("DOCUMENT_SOURCE_TIME_INVALID", "原文件创建时间格式无效", 422)
        try:
            datetime.fromisoformat(timestamp.replace("Z", "+00:00"))
        except ValueError as error:
            raise AppError("DOCUMENT_SOURCE_TIME_INVALID", "原文件创建时间格式无效", 422) from error
        return timestamp

    def list_documents(
        self, include_deleted: bool = False, course_id: int | None = None
    ) -> list[dict]:
        deleted_clause = "" if include_deleted else "AND deleted_at IS NULL"
        with self.database.connect() as connection:
            rows = connection.execute(
                f"""SELECT * FROM documents
                WHERE course_id = ? {deleted_clause}
                ORDER BY id DESC""",
                (self._course_id() if course_id is None else int(course_id),),
            ).fetchall()
        return [self._serialize(row) for row in rows]

    def get_document(self, document_id: int, include_deleted: bool = False) -> dict:
        deleted_clause = "" if include_deleted else "AND deleted_at IS NULL"
        with self.database.connect() as connection:
            row = connection.execute(
                f"""SELECT * FROM documents
                WHERE id = ? AND course_id = ? {deleted_clause}""",
                (document_id, self._course_id()),
            ).fetchone()
        if not row:
            raise AppError("DOCUMENT_NOT_FOUND", "文档不存在", 404)
        return self._serialize(row)

    def original_file(self, document_id: int) -> tuple[Path, str, str]:
        document = self.get_document(document_id)
        root = self.data_dir.resolve()
        path = (root / str(document["stored_path"])).resolve()
        try:
            path.relative_to(root)
        except ValueError as error:
            raise AppError("DOCUMENT_FILE_INVALID", "资料文件路径无效", 500) from error
        if not path.is_file():
            raise AppError("DOCUMENT_FILE_NOT_FOUND", "资料原文件不存在", 404)
        return path, str(document["media_type"]), str(document["filename"])

    def export_document(self, document_id: int, format_name: str) -> tuple[bytes, str, str]:
        if format_name not in {"source", "pdf"}:
            raise AppError("UNSUPPORTED_EXPORT", "不支持的资料导出格式", 422)
        content = self.get_content(document_id)
        document = content["document"]
        blocks = content["blocks"]
        safe_title = self._safe_export_filename(str(document["title"]))
        source_path, source_media_type, source_filename = self.original_file(document_id)
        source_suffix = Path(source_filename).suffix.lower()

        if format_name == "source":
            if document["format"] == "csv":
                cells = (blocks[0].get("data") or {}).get("cells") if blocks else []
                cells = cells if isinstance(cells, list) else []
                max_row = max((int(cell.get("row") or 0) for cell in cells if isinstance(cell, dict)), default=0)
                max_column = max((int(cell.get("column") or 0) for cell in cells if isinstance(cell, dict)), default=0)
                values = {
                    (int(cell.get("row") or 0), int(cell.get("column") or 0)): str(cell.get("value") or "")
                    for cell in cells if isinstance(cell, dict)
                }
                output = StringIO(newline="")
                delimiter = "\t" if source_suffix == ".tsv" else ","
                writer = csv.writer(output, delimiter=delimiter)
                for row in range(1, max_row + 1):
                    writer.writerow([values.get((row, column), "") for column in range(1, max_column + 1)])
                suffix = ".tsv" if source_suffix == ".tsv" else ".csv"
                media_type = "text/tab-separated-values; charset=utf-8" if suffix == ".tsv" else "text/csv; charset=utf-8"
                return output.getvalue().encode("utf-8-sig"), f"{safe_title}{suffix}", media_type
            if document["format"] in {"markdown", "text"}:
                if document["format"] == "markdown":
                    suffix = source_suffix if source_suffix in {".md", ".markdown"} else ".md"
                    media_type = source_media_type or "text/markdown; charset=utf-8"
                else:
                    suffix = source_suffix if source_suffix in TEXT_SUFFIXES else ".txt"
                    media_type = source_media_type or "text/plain; charset=utf-8"
                body = "\n\n".join(str(block.get("text") or "").rstrip() for block in blocks).rstrip() + "\n"
                return body.encode("utf-8"), f"{safe_title}{suffix}", media_type
            if document["format"] == "docx":
                return (
                    self._render_saved_docx(blocks),
                    f"{safe_title}.docx",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            return source_path.read_bytes(), f"{safe_title}{source_suffix}", source_media_type

        if document["format"] == "pdf":
            return source_path.read_bytes(), f"{safe_title}.pdf", "application/pdf"
        body = "\n\n".join(str(block.get("text") or "").strip() for block in blocks if str(block.get("text") or "").strip())
        return self._render_text_pdf(safe_title, body), f"{safe_title}.pdf", "application/pdf"

    @staticmethod
    def _safe_export_filename(value: str) -> str:
        cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "-", value).strip(" .")
        return (cleaned or "StudyPilot-资料")[:120]

    @staticmethod
    def _render_saved_docx(blocks: list[dict]) -> bytes:
        exported = Document()
        for block in blocks:
            text = str(block.get("text") or "").strip()
            if not text:
                continue
            data = block.get("data") or {}
            rows = data.get("rows") if block.get("block_type") == "table" else None
            if isinstance(rows, list) and rows:
                width = max((len(row) for row in rows if isinstance(row, list)), default=1)
                table = exported.add_table(rows=len(rows), cols=max(1, width))
                for row_index, row in enumerate(rows):
                    if not isinstance(row, list):
                        continue
                    for column_index, value in enumerate(row[:width]):
                        table.cell(row_index, column_index).text = str(value)
                continue
            style = str(data.get("style") or "")
            heading = re.search(r"heading\s*(\d+)", style, flags=re.IGNORECASE)
            if heading:
                exported.add_heading(text, level=max(1, min(6, int(heading.group(1)))))
            else:
                exported.add_paragraph(text)
        stream = BytesIO()
        exported.save(stream)
        return stream.getvalue()

    @staticmethod
    def _render_text_pdf(title: str, body: str) -> bytes:
        stream = BytesIO()
        page_width, page_height = A4
        pdf = canvas.Canvas(stream, pagesize=A4)
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            font_name = "STSong-Light"
        except Exception:
            font_name = "Helvetica"
        pdf.setTitle(title)
        margin = 52
        y = page_height - margin
        pdf.setFont(font_name, 18)
        pdf.drawString(margin, y, title)
        y -= 32
        pdf.setFont(font_name, 11)
        max_width = page_width - margin * 2
        for paragraph in (body or "暂无可导出的正文").splitlines():
            lines = DocumentService._wrap_pdf_line(paragraph, font_name, 11, max_width) or [""]
            for line in lines:
                if y < margin:
                    pdf.showPage()
                    pdf.setFont(font_name, 11)
                    y = page_height - margin
                pdf.drawString(margin, y, line)
                y -= 17
            y -= 5
        pdf.save()
        return stream.getvalue()

    @staticmethod
    def _wrap_pdf_line(value: str, font_name: str, font_size: int, max_width: float) -> list[str]:
        lines: list[str] = []
        current = ""
        for character in value:
            candidate = current + character
            if current and pdfmetrics.stringWidth(candidate, font_name, font_size) > max_width:
                lines.append(current)
                current = character
            else:
                current = candidate
        if current:
            lines.append(current)
        return lines

    def update_document(self, document_id: int, values: dict) -> dict:
        current = self.get_document(document_id)
        title = str(values.get("title") or current["title"]).strip()
        metadata = dict(current.get("metadata") or {})
        if values.get("favorite") is not None:
            metadata["favorite"] = bool(values["favorite"])
        if values.get("pinned") is not None:
            metadata["pinned"] = bool(values["pinned"])
        with self.database.connect() as connection:
            connection.execute(
                """UPDATE documents
                SET title = ?, metadata_json = ?, updated_at = CURRENT_TIMESTAMP
                WHERE id = ? AND course_id = ? AND deleted_at IS NULL""",
                (title, json.dumps(metadata, ensure_ascii=False), document_id, self._course_id()),
            )
            row = connection.execute(
                "SELECT * FROM documents WHERE id = ?", (document_id,)
            ).fetchone()
        return self._serialize(row)

    def trash_document(self, document_id: int) -> None:
        self.get_document(document_id)
        with self.database.connect() as connection:
            connection.execute(
                """UPDATE documents
                SET deleted_at = CURRENT_TIMESTAMP, updated_at = CURRENT_TIMESTAMP
                WHERE id = ? AND course_id = ?""",
                (document_id, self._course_id()),
            )

    def restore_document(self, document_id: int) -> dict:
        self.get_document(document_id, include_deleted=True)
        with self.database.connect() as connection:
            connection.execute(
                """UPDATE documents
                SET deleted_at = NULL, updated_at = CURRENT_TIMESTAMP
                WHERE id = ? AND course_id = ?""",
                (document_id, self._course_id()),
            )
            row = connection.execute(
                "SELECT * FROM documents WHERE id = ?", (document_id,)
            ).fetchone()
        return self._serialize(row)

    def get_content(self, document_id: int) -> dict:
        document = self.get_document(document_id)
        with self.database.connect() as connection:
            rows = connection.execute(
                "SELECT * FROM document_blocks WHERE document_id = ? ORDER BY ordinal, id",
                (document_id,),
            ).fetchall()
        blocks = [self._serialize_block(row) for row in rows]
        return {
            "document": document,
            "blocks": self._enrich_heading_hierarchy(document, blocks),
        }

    @staticmethod
    def _enrich_heading_hierarchy(document: dict, blocks: list[dict]) -> list[dict]:
        """Backfill outline metadata for documents imported by an older release."""
        if document.get("format") not in {"markdown", "docx"}:
            return blocks
        heading_stack: dict[int, dict] = {}
        for block in blocks:
            data = dict(block.get("data") or {})
            level = int(data.get("heading_level") or 0)
            if not level and document.get("format") == "markdown":
                heading = re.match(r"^\s{0,3}(#{1,6})\s+(.+?)\s*(?:\n|$)", str(block.get("text") or ""))
                if heading:
                    level = len(heading.group(1))
                    data.setdefault("title", heading.group(2).strip())
            elif not level and document.get("format") == "docx":
                style = str(data.get("style") or "")
                style_match = re.search(r"heading\s*(\d+)", style, flags=re.IGNORECASE)
                if style_match:
                    level = max(1, min(6, int(style_match.group(1))))

            if not level:
                block["data"] = data
                continue

            heading_stack = {
                stack_level: stack_block
                for stack_level, stack_block in heading_stack.items()
                if stack_level < level
            }
            parent = next(
                (
                    heading_stack[parent_level]
                    for parent_level in range(level - 1, 0, -1)
                    if parent_level in heading_stack
                ),
                None,
            )
            title = str(data.get("title") or block.get("text") or f"章节 {block.get('ordinal', 0) + 1}").strip()
            data.update(
                {
                    "heading_level": level,
                    "parent_block_key": parent["block_key"] if parent else "",
                    "heading_path": [
                        str(heading_stack[stack_level]["data"].get("title") or "")
                        for stack_level in sorted(heading_stack)
                    ] + [title],
                }
            )
            block["data"] = data
            heading_stack[level] = block
        return blocks

    def add_revision(self, document_id: int, values: dict) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            block = self._block_row(connection, document_id, values["block_key"])
            before = values.get("before") or {}
            after = values["after"]
            if "text" in before and before["text"] != block["text"]:
                raise AppError(
                    "DOCUMENT_REVISION_CONFLICT",
                    "资料内容已在其他位置更新，请刷新后重试",
                    409,
                )
            connection.execute(
                "DELETE FROM document_revisions WHERE document_id = ? AND is_applied = 0",
                (document_id,),
            )
            next_revision = int(
                connection.execute(
                    """SELECT COALESCE(MAX(revision), 0) + 1
                    FROM document_revisions WHERE document_id = ? AND block_key = ?""",
                    (document_id, values["block_key"]),
                ).fetchone()[0]
            )
            cursor = connection.execute(
                """INSERT INTO document_revisions(
                    document_id, block_key, before_json, after_json, revision
                ) VALUES (?, ?, ?, ?, ?)""",
                (
                    document_id,
                    values["block_key"],
                    json.dumps(before, ensure_ascii=False),
                    json.dumps(after, ensure_ascii=False),
                    next_revision,
                ),
            )
            self._apply_block_state(connection, document_id, block, after)
            revision_row = connection.execute(
                "SELECT * FROM document_revisions WHERE id = ?", (cursor.lastrowid,)
            ).fetchone()
            updated_block = connection.execute(
                "SELECT * FROM document_blocks WHERE id = ?", (block["id"],)
            ).fetchone()
        return {
            "revision": self._serialize_revision(revision_row),
            "block": self._serialize_block(updated_block),
            "history": self.revision_state(document_id),
        }

    def revision_state(self, document_id: int) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            row = connection.execute(
                """SELECT
                    EXISTS(SELECT 1 FROM document_revisions WHERE document_id = ? AND is_applied = 1) AS can_undo,
                    EXISTS(SELECT 1 FROM document_revisions WHERE document_id = ? AND is_applied = 0) AS can_redo""",
                (document_id, document_id),
            ).fetchone()
        return {"can_undo": bool(row["can_undo"]), "can_redo": bool(row["can_redo"])}

    def undo_revision(self, document_id: int) -> dict:
        return self._move_revision(document_id, undo=True)

    def redo_revision(self, document_id: int) -> dict:
        return self._move_revision(document_id, undo=False)

    def _move_revision(self, document_id: int, *, undo: bool) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            revision = connection.execute(
                f"""SELECT * FROM document_revisions
                WHERE document_id = ? AND is_applied = ?
                ORDER BY id {'DESC' if undo else 'ASC'} LIMIT 1""",
                (document_id, 1 if undo else 0),
            ).fetchone()
            if not revision:
                raise AppError(
                    "DOCUMENT_HISTORY_EMPTY",
                    "没有可以撤销的修改" if undo else "没有可以重做的修改",
                    409,
                )
            block = self._block_row(connection, document_id, revision["block_key"])
            state = json.loads(revision["before_json"] if undo else revision["after_json"])
            self._apply_block_state(connection, document_id, block, state)
            connection.execute(
                "UPDATE document_revisions SET is_applied = ? WHERE id = ?",
                (0 if undo else 1, revision["id"]),
            )
            updated_block = connection.execute(
                "SELECT * FROM document_blocks WHERE id = ?", (block["id"],)
            ).fetchone()
            state_row = connection.execute(
                """SELECT
                    EXISTS(SELECT 1 FROM document_revisions WHERE document_id = ? AND is_applied = 1) AS can_undo,
                    EXISTS(SELECT 1 FROM document_revisions WHERE document_id = ? AND is_applied = 0) AS can_redo""",
                (document_id, document_id),
            ).fetchone()
        return {
            "block": self._serialize_block(updated_block),
            "history": {"can_undo": bool(state_row["can_undo"]), "can_redo": bool(state_row["can_redo"])},
        }

    @staticmethod
    def _apply_block_state(connection, document_id: int, block, state: dict) -> None:
        text = str(state.get("text", block["text"]))
        data_json = block["data_json"] if "data" not in state else json.dumps(state["data"], ensure_ascii=False)
        connection.execute(
            """UPDATE document_blocks
            SET text = ?, data_json = ?, updated_at = CURRENT_TIMESTAMP
            WHERE id = ?""",
            (text, data_json, block["id"]),
        )
        body = "\n\n".join(
            row["text"].strip()
            for row in connection.execute(
                """SELECT text FROM document_blocks
                WHERE document_id = ? ORDER BY ordinal, id""",
                (document_id,),
            ).fetchall()
            if row["text"].strip()
        )
        connection.execute(
            "UPDATE documents SET body = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
            (body, document_id),
        )

    def list_annotations(self, document_id: int) -> list[dict]:
        self.get_document(document_id)
        with self.database.connect() as connection:
            rows = connection.execute(
                """SELECT * FROM document_annotations
                WHERE document_id = ? ORDER BY id""",
                (document_id,),
            ).fetchall()
        return [self._serialize_annotation(row) for row in rows]

    def add_annotation(self, document_id: int, values: dict) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            self._block_row(connection, document_id, values["block_key"])
            cursor = connection.execute(
                """INSERT INTO document_annotations(
                    document_id, block_key, kind, locator_json, quote, note,
                    color, geometry_json
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
                (
                    document_id,
                    values["block_key"],
                    values["kind"],
                    json.dumps(values.get("locator") or {}, ensure_ascii=False),
                    values.get("quote", ""),
                    values.get("note", ""),
                    values.get("color", "yellow"),
                    json.dumps(values.get("geometry") or {}, ensure_ascii=False),
                ),
            )
            row = connection.execute(
                "SELECT * FROM document_annotations WHERE id = ?", (cursor.lastrowid,)
            ).fetchone()
        return self._serialize_annotation(row)

    def update_annotation(self, document_id: int, annotation_id: int, values: dict) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            current = connection.execute(
                """SELECT * FROM document_annotations
                WHERE id = ? AND document_id = ?""",
                (annotation_id, document_id),
            ).fetchone()
            if not current:
                raise AppError("ANNOTATION_NOT_FOUND", "批注不存在", 404)
            fields: dict[str, Any] = {}
            for name in ("note", "color"):
                if name in values:
                    fields[name] = values[name]
            if "geometry" in values:
                fields["geometry_json"] = json.dumps(values["geometry"] or {}, ensure_ascii=False)
            if fields:
                assignments = ", ".join(f"{name} = ?" for name in fields)
                connection.execute(
                    f"""UPDATE document_annotations SET {assignments},
                    revision = revision + 1, updated_at = CURRENT_TIMESTAMP
                    WHERE id = ? AND document_id = ?""",
                    [*fields.values(), annotation_id, document_id],
                )
            row = connection.execute(
                "SELECT * FROM document_annotations WHERE id = ?", (annotation_id,)
            ).fetchone()
        return self._serialize_annotation(row)

    def delete_annotation(self, document_id: int, annotation_id: int) -> None:
        self.get_document(document_id)
        with self.database.connect() as connection:
            cursor = connection.execute(
                """DELETE FROM document_annotations
                WHERE id = ? AND document_id = ?""",
                (annotation_id, document_id),
            )
            if not cursor.rowcount:
                raise AppError("ANNOTATION_NOT_FOUND", "批注不存在", 404)

    def add_highlight(self, document_id: int, values: dict) -> dict:
        self.get_document(document_id)
        with self.database.connect() as connection:
            cursor = connection.execute(
                """INSERT INTO document_highlights(
                    document_id, quote, note, start_offset, end_offset
                ) VALUES (?, ?, ?, ?, ?)""",
                (
                    document_id,
                    values["quote"],
                    values.get("note", ""),
                    values.get("start_offset"),
                    values.get("end_offset"),
                ),
            )
            row = connection.execute(
                "SELECT * FROM document_highlights WHERE id = ?", (cursor.lastrowid,)
            ).fetchone()
        return as_dict(row)

    def search(self, query: str, limit: int = 20) -> list[dict]:
        terms = re.findall(r"[\w\u3400-\u9fff]+", query, flags=re.UNICODE)
        if not terms:
            return []
        fts_query = " AND ".join(f'"{term}"' for term in terms[:10])
        with self.database.connect() as connection:
            rows = connection.execute(
                """SELECT d.id AS document_id, d.title, d.filename,
                    snippet(document_fts, 1, '<mark>', '</mark>', '…', 24) AS snippet,
                    bm25(document_fts) AS score
                FROM document_fts JOIN documents d ON d.id = document_fts.rowid
                WHERE document_fts MATCH ? AND d.course_id = ? AND d.deleted_at IS NULL
                ORDER BY score LIMIT ?""",
                (fts_query, self._course_id(), limit),
            ).fetchall()
        return [as_dict(row) for row in rows]

    @staticmethod
    def _serialize(row) -> dict:
        item = as_dict(row)
        item["metadata"] = json.loads(item.pop("metadata_json"))
        item["structure"] = json.loads(item.pop("structure_json", "{}"))
        return item

    @staticmethod
    def _serialize_block(row) -> dict:
        item = as_dict(row)
        item["locator"] = json.loads(item.pop("locator_json"))
        item["data"] = json.loads(item.pop("data_json"))
        return item

    @staticmethod
    def _serialize_revision(row) -> dict:
        item = as_dict(row)
        item["before"] = json.loads(item.pop("before_json"))
        item["after"] = json.loads(item.pop("after_json"))
        return item

    @staticmethod
    def _serialize_annotation(row) -> dict:
        item = as_dict(row)
        item["locator"] = json.loads(item.pop("locator_json"))
        item["geometry"] = json.loads(item.pop("geometry_json"))
        return item

    @staticmethod
    def _block_row(connection, document_id: int, block_key: str):
        row = connection.execute(
            """SELECT * FROM document_blocks
            WHERE document_id = ? AND block_key = ?""",
            (document_id, block_key),
        ).fetchone()
        if not row:
            raise AppError("DOCUMENT_BLOCK_NOT_FOUND", "资料内容块不存在", 404)
        return row
